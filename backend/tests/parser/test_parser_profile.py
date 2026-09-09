"""Tests for document profiling (_profile.py)."""

import os
import tempfile
from pathlib import Path

import pytest

from src.core import constants as constants_module

os.environ["API_KEY"] = ""
os.environ["DATA_DIR"] = tempfile.mkdtemp()

constants_module.DATA_DIR = Path(os.environ["DATA_DIR"])
constants_module.DATABASE_PATH = constants_module.DATA_DIR / "test.db"
constants_module.CHROMA_PATH = constants_module.DATA_DIR / "chroma"
constants_module.UPLOAD_DIR = constants_module.DATA_DIR / "uploads"

from src.services.parser._profile import DocumentProfile, profile_document  # noqa: E402


class TestDocumentProfile:
    """Test DocumentProfile dataclass immutability."""

    def test_frozen_dataclass(self):
        profile = DocumentProfile(
            suffix=".pdf",
            page_count=10,
            avg_chars_per_page=250.0,
            image_ratio=0.2,
            has_text_layer=True,
            is_likely_scanned=False,
            size_bytes=1024,
        )
        with pytest.raises(AttributeError):
            profile.page_count = 20  # type: ignore[misc]


class TestProfileImageFile:
    """Test profiling of single image files."""

    def test_image_profile(self, tmp_path):
        img = tmp_path / "test.png"
        img.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
        profile = profile_document(img)

        assert profile.suffix == ".png"
        assert profile.page_count == 1
        assert profile.image_ratio == 1.0
        assert profile.has_text_layer is False
        assert profile.is_likely_scanned is False

    def test_jpeg_profile(self, tmp_path):
        img = tmp_path / "photo.jpg"
        img.write_bytes(b"\xff\xd8\xff\xe0" + b"\x00" * 100)
        profile = profile_document(img)
        assert profile.suffix == ".jpg"
        assert profile.image_ratio == 1.0


class TestProfilePDF:
    """Test profiling of PDF files via PyMuPDF."""

    @pytest.fixture
    def text_pdf(self, tmp_path):
        """Create a simple text-only PDF."""
        import pymupdf as fitz

        path = tmp_path / "text_only.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "A" * 500)  # Dense text
        doc.save(str(path))
        doc.close()
        return path

    @pytest.fixture
    def scanned_pdf(self, tmp_path):
        """Create a scanned (image-only) PDF."""
        import pymupdf as fitz

        path = tmp_path / "scanned.pdf"
        doc = fitz.open()
        page = doc.new_page()
        # Insert an image so the page has image content but no text
        pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 200, 200))
        pix.set_rect(pix.irect, (255, 255, 255))
        page.insert_image(page.rect, pixmap=pix)
        doc.save(str(path))
        doc.close()
        return path

    def test_text_pdf_has_text_layer(self, text_pdf):
        profile = profile_document(text_pdf)
        assert profile.suffix == ".pdf"
        assert profile.page_count == 1
        assert profile.avg_chars_per_page > 0
        assert profile.has_text_layer is True
        assert profile.is_likely_scanned is False

    def test_scanned_pdf_detected(self, scanned_pdf):
        profile = profile_document(scanned_pdf)
        assert profile.suffix == ".pdf"
        assert profile.page_count == 1
        assert profile.is_likely_scanned is True
        assert profile.has_text_layer is False

    def test_nonexistent_pdf_returns_unknown(self, tmp_path):
        fake = tmp_path / "nonexistent.pdf"
        profile = profile_document(fake)
        assert profile.suffix == ".pdf"
        assert profile.page_count == 0
        assert profile.avg_chars_per_page == 0.0


class TestProfilePPTX:
    """Test profiling of PPTX files via python-pptx."""

    @pytest.fixture
    def text_pptx(self, tmp_path):
        """Create a PPTX with text shapes."""
        from pptx import Presentation

        path = tmp_path / "text_slides.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Title"
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = "Body " * 100
        prs.save(str(path))
        return path

    def test_text_pptx_profile(self, text_pptx):
        profile = profile_document(text_pptx)
        assert profile.suffix == ".pptx"
        assert profile.page_count == 1
        assert profile.avg_chars_per_page > 0
        assert profile.has_text_layer is True
        assert profile.image_ratio < 0.5


class TestProfileOtherFormats:
    """Test profiling of non-PDF/PPTX formats."""

    def test_docx_profile(self, tmp_path):
        docx = tmp_path / "test.docx"
        docx.write_bytes(b"PK" + b"\x00" * 100)
        profile = profile_document(docx)
        assert profile.suffix == ".docx"
        assert profile.has_text_layer is True

    def test_unknown_format_returns_generic(self, tmp_path):
        f = tmp_path / "test.xyz"
        f.write_bytes(b"data")
        profile = profile_document(f)
        assert profile.suffix == ".xyz"
