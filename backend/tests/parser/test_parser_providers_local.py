"""Tests for local fast parser (providers/local.py)."""

import os
import tempfile
from pathlib import Path

import pytest

from src.core import constants as constants_module
from src.core.config import settings

os.environ["API_KEY"] = ""
os.environ["DATA_DIR"] = tempfile.mkdtemp()

constants_module.DATA_DIR = Path(os.environ["DATA_DIR"])
constants_module.DATABASE_PATH = constants_module.DATA_DIR / "test.db"
constants_module.CHROMA_PATH = constants_module.DATA_DIR / "chroma"
constants_module.UPLOAD_DIR = constants_module.DATA_DIR / "uploads"

from src.services.parser.providers.local import LocalFastParser  # noqa: E402
from src.services.parser.types import ParsedDocument  # noqa: E402


@pytest.fixture
def parser():
    return LocalFastParser()


@pytest.fixture
def text_pdf(tmp_path):
    """Create a simple text-only PDF."""
    import fitz

    path = tmp_path / "text.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello World from PDF")
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def text_pptx(tmp_path):
    """Create a simple PPTX with text."""
    from pptx import Presentation

    path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    prs.save(str(path))
    return path


class TestLocalPDF:
    """Test local PDF parsing."""

    def test_parse_pdf(self, parser, text_pdf):
        result = parser.parse(text_pdf)
        assert isinstance(result, ParsedDocument)
        assert result.file_type == "pdf"
        assert result.total_pages == 1
        assert "Hello World" in result.to_text()

    def test_pdf_metadata(self, parser, text_pdf):
        result = parser.parse(text_pdf)
        assert result.metadata["parser"] == "local"
        assert result.metadata["has_ocr"] is False

    def test_multi_page_pdf(self, tmp_path, parser):
        import fitz

        path = tmp_path / "multi.pdf"
        doc = fitz.open()
        for i in range(3):
            page = doc.new_page()
            page.insert_text((72, 72), f"Page {i + 1}")
        doc.save(str(path))
        doc.close()

        result = parser.parse(path)
        assert result.total_pages == 3

    def test_empty_pdf(self, tmp_path, parser):
        import fitz

        path = tmp_path / "empty.pdf"
        doc = fitz.open()
        doc.new_page()
        doc.save(str(path))
        doc.close()

        result = parser.parse(path)
        assert result.total_pages == 1
        assert result.to_text() == ""

    def test_empty_pdf_renders_page_fallback_image(self, tmp_path, parser, monkeypatch):
        import fitz

        monkeypatch.setattr(settings, "RAG_INDEX_IMAGE_CAPTIONS", True)
        path = tmp_path / "empty_no_text.pdf"
        doc = fitz.open()
        doc.new_page(width=595, height=842)
        doc.save(str(path))
        doc.close()

        result = parser.parse(path)
        assert result.total_pages == 1
        assert result.pages[0].images is not None
        assert len(result.pages[0].images) == 1
        assert result.pages[0].images[0]["name"].endswith(".png")


class TestLocalPPTX:
    """Test local PPTX parsing."""

    def test_parse_pptx(self, parser, text_pptx):
        result = parser.parse(text_pptx)
        assert isinstance(result, ParsedDocument)
        assert result.file_type == "pptx"
        assert "Slide Title" in result.to_text()

    def test_pptx_metadata(self, parser, text_pptx):
        result = parser.parse(text_pptx)
        assert result.metadata["parser"] == "local"
        assert result.metadata["has_ocr"] is False

    def test_pptx_with_table(self, tmp_path, parser):
        from pptx import Presentation
        from pptx.util import Inches

        path = tmp_path / "table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        rows, cols = 2, 2
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
        table_shape.table.cell(0, 0).text = "A"
        table_shape.table.cell(0, 1).text = "B"
        table_shape.table.cell(1, 0).text = "C"
        table_shape.table.cell(1, 1).text = "D"
        prs.save(str(path))

        result = parser.parse(path)
        assert "[TABLE]" in result.to_text()
        assert result.pages[0].tables is not None


class TestLocalUnsupportedFormat:
    """Test local parser with unknown formats."""

    def test_generic_fallback(self, tmp_path, parser):
        f = tmp_path / "test.txt"
        f.write_text("plain text content")
        result = parser.parse(f)
        assert isinstance(result, ParsedDocument)
