"""Generate synthetic benchmark fixtures.

Run this script to create the fixture files. Commit the generated files so
CI and other developers don't need to run the generator.

Usage:
    cd backend/tests/fixtures/benchmark
    uv run python generate.py
"""

import math
import struct
import wave
from pathlib import Path

from pptx import Presentation  # type: ignore[import-not-found]
from pptx.util import Inches  # type: ignore[import-not-found]

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
except ImportError as exc:
    raise ImportError(
        "reportlab is required to generate PDF fixtures. Install it with: pip install reportlab"
    ) from exc

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError as exc:
    raise ImportError("Pillow is required to generate scanned PDF fixtures.") from exc

FIXTURE_DIR = Path(__file__).parent


SAMPLE_TEXT = """
Q1 Planning Meeting Notes

Attendees: Alice, Bob, Carol, Dave
Date: January 15, 2026

Agenda:
1. Review Q4 results
2. Discuss Q1 priorities
3. Identify blockers and risks

Discussion:
Alice opened the meeting by highlighting the strong Q4 performance. Revenue was up 12% quarter-over-quarter, and customer retention reached an all-time high of 94%.
Bob presented the engineering roadmap for Q1. The top priorities are:
- Launch the new mobile app by February 15
- Migrate the legacy database to PostgreSQL by March 1
- Implement the new analytics pipeline

Blockers raised:
- The mobile app is waiting on final UI designs from the design team. ETA is January 22.
- Database migration is blocked by the need for additional cloud storage budget approval.
- The analytics pipeline depends on a third-party API that has not yet provided production access.

Action items:
- Alice to follow up with the design team for UI deliverables.
- Bob to submit a budget request for cloud storage by January 18.
- Carol to contact the third-party API provider and request production credentials.

Next meeting: January 29, 2026
""".strip()


def generate_sample_pdf() -> None:
    """Create a text-heavy 2-page PDF (marker path)."""
    path = FIXTURE_DIR / "sample.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    text_object = c.beginText(40, height - 40)
    text_object.setFont("Helvetica", 11)

    lines = SAMPLE_TEXT.splitlines()
    for line in lines:
        if text_object.getY() < 60:
            c.drawText(text_object)
            c.showPage()
            text_object = c.beginText(40, height - 40)
            text_object.setFont("Helvetica", 11)
        text_object.textLine(line)

    c.drawText(text_object)
    c.save()
    print(f"Generated {path} ({path.stat().st_size} bytes)")


def generate_scanned_pdf() -> None:
    """Create an image-only 1-page PDF that forces OCR fallback."""
    path = FIXTURE_DIR / "scanned.pdf"
    img = Image.new("RGB", (612, 792), color="white")
    draw = ImageDraw.Draw(img)

    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
    except Exception:
        font = ImageFont.load_default()

    text = "Scanned Document\n\nThis is a test page for OCR fallback benchmarking.\n"
    draw.multiline_text((50, 50), text, fill="black", font=font)
    img.save(path, "PDF", resolution=100.0)
    print(f"Generated {path} ({path.stat().st_size} bytes)")


def generate_sample_pptx() -> None:
    """Create a 3-slide PPTX."""
    path = FIXTURE_DIR / "sample.pptx"
    prs = Presentation()

    slide_layout = prs.slide_layouts[6]  # blank layout

    slide1 = prs.slides.add_slide(slide_layout)
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.text_frame.text = "Q1 Planning Meeting"
    body_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    body_box.text_frame.text = (
        "Attendees: Alice, Bob, Carol, Dave\n"
        "Date: January 15, 2026\n"
        "Revenue up 12% QoQ, retention at 94%"
    )

    slide2 = prs.slides.add_slide(slide_layout)
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.text_frame.text = "Engineering Roadmap"
    body_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    body_box.text_frame.text = (
        "1. Launch mobile app by Feb 15\n"
        "2. Migrate to PostgreSQL by Mar 1\n"
        "3. Implement analytics pipeline"
    )

    slide3 = prs.slides.add_slide(slide_layout)
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.text_frame.text = "Blockers & Action Items"
    body_box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    body_box.text_frame.text = (
        "Blockers:\n"
        "- Mobile app waiting on UI designs\n"
        "- DB migration blocked by budget\n"
        "- Analytics pipeline needs API access\n\n"
        "Actions:\n"
        "- Alice: follow up with design\n"
        "- Bob: submit budget request\n"
        "- Carol: request API credentials"
    )

    prs.save(path)
    print(f"Generated {path} ({path.stat().st_size} bytes)")


def generate_sample_wav() -> None:
    """Create a 15-second mono 16 kHz sine-sweep WAV."""
    path = FIXTURE_DIR / "sample.wav"
    duration = 15.0
    sample_rate = 16000
    amplitude = 0.3

    n_samples = int(duration * sample_rate)
    with wave.open(str(path), "w") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(sample_rate)
        for i in range(n_samples):
            t = i / sample_rate
            freq = 440 + (880 - 440) * (t / duration)
            sample = amplitude * math.sin(2 * math.pi * freq * t)
            w.writeframes(struct.pack("<h", int(sample * 32767)))

    print(f"Generated {path} ({path.stat().st_size} bytes)")


if __name__ == "__main__":
    FIXTURE_DIR.mkdir(parents=True, exist_ok=True)
    generate_sample_pdf()
    generate_scanned_pdf()
    generate_sample_pptx()
    generate_sample_wav()
    print("All fixtures generated successfully.")
