"""PaddleOCR layout-parsing API client — aistudio-app endpoint.

Primarily used for single-image layout extraction. For multi-page PDFs,
renders each page via PyMuPDF and calls the API per page (bounded parallelism).
"""

import asyncio
import base64
import logging
from pathlib import Path
from typing import Any

from src.core.config import settings

from .._errors import ParserProviderError
from .._http import get_parser_http_client
from .._profile import DocumentProfile
from ..types import PageContent, ParsedDocument

logger = logging.getLogger(__name__)

_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"})

# Max concurrent page requests for multi-page PDFs
_MAX_CONCURRENT_PAGES = 4


class PaddleOCRAPIParser:
    """HTTP client for the aistudio-app PaddleOCR layout-parsing endpoint."""

    async def parse(self, file_path: Path, profile: DocumentProfile) -> ParsedDocument:
        """Parse a document via the PaddleOCR layout-parsing API.

        Raises ``ParserProviderError`` on any failure.
        """
        api_key = settings.PADDLEOCR_API_KEY.get_secret_value()
        if not api_key:
            raise ParserProviderError(provider="paddle", retryable=False)

        base_url = settings.PADDLEOCR_BASE_URL.rstrip("/")
        if not base_url:
            raise ParserProviderError(
                provider="paddle",
                retryable=False,
                cause=Exception("PADDLEOCR_BASE_URL not configured"),
            )

        client = get_parser_http_client()
        suffix = file_path.suffix.lower()

        try:
            if suffix in _IMAGE_EXTS:
                return await self._parse_image(client, base_url, api_key, file_path)
            if suffix == ".pdf":
                return await self._parse_pdf(client, base_url, api_key, file_path)
            if suffix == ".pptx":
                return await self._parse_pptx(client, base_url, api_key, file_path)

            raise ParserProviderError(
                provider="paddle",
                retryable=False,
                cause=ValueError(f"Unsupported format for PaddleOCR: {suffix}"),
            )
        except ParserProviderError:
            raise
        except Exception as exc:
            raise ParserProviderError(provider="paddle", retryable=True, cause=exc) from exc

    async def _call_layout_api(
        self,
        client,
        base_url: str,
        api_key: str,
        image_bytes: bytes,
    ) -> list[dict[str, Any]]:
        """Call the layout-parsing API with a single image.

        Returns the list of layout blocks.
        """
        from .._http import get_parser_http_client

        img_b64 = base64.b64encode(image_bytes).decode("ascii")
        payload = {"file": img_b64}
        headers = {"Authorization": f"token {api_key}"}

        resp = None
        for attempt in range(2):
            try:
                resp = await client.post(
                    base_url,
                    json=payload,
                    headers=headers,
                )
                break
            except RuntimeError as exc:
                if attempt == 0 and "Event loop is closed" in str(exc):
                    logger.warning("PaddleOCR client hit closed loop, recreating and retrying")
                    client = get_parser_http_client()
                    continue
                raise
        if resp is None:
            raise ParserProviderError(
                provider="paddle",
                retryable=True,
                cause=RuntimeError("PaddleOCR API call failed without response"),
            )

        if resp.status_code in (401, 403):
            logger.error(
                "PaddleOCR API key invalid or revoked (HTTP %d). "
                "Please check your PADDLEOCR_API_KEY configuration.",
                resp.status_code,
            )
            raise ParserProviderError(
                provider="paddle",
                retryable=False,
                cause=Exception(
                    f"API key invalid/revoked (HTTP {resp.status_code}) — check PADDLEOCR_API_KEY"
                ),
            )
        if resp.status_code >= 500:
            raise ParserProviderError(
                provider="paddle",
                retryable=True,
                cause=Exception(f"Server error {resp.status_code}"),
            )
        if resp.status_code >= 400:
            logger.debug(
                "Provider paddle HTTP %d response: %s",
                resp.status_code,
                resp.text[:500],
            )
            raise ParserProviderError(
                provider="paddle",
                retryable=False,
                cause=Exception(f"HTTP {resp.status_code} from paddle"),
            )

        body = resp.json()

        # The response may contain layout blocks in various structures
        # depending on the PaddleOCR Serving version.
        # aistudio-app format:
        # {"result":{"layoutParsingResults":[{"prunedResult":{"parsing_res_list":[...]}}]}}
        if isinstance(body, list):
            return body
        if isinstance(body, dict):
            # Try nested aistudio-app format first
            result = body.get("result")
            if isinstance(result, dict):
                lpr = result.get("layoutParsingResults", [])
                if isinstance(lpr, list) and lpr:
                    pruned = lpr[0].get("prunedResult", {})
                    blocks = pruned.get("parsing_res_list", [])
                    if isinstance(blocks, list):
                        return blocks
            # Legacy / alternative formats
            return body.get("layout", body.get("blocks", []))
        return []

    def _blocks_to_text(self, blocks: list[dict[str, Any]]) -> tuple[str, list[list[list[str]]]]:
        """Convert layout blocks to text + tables.

        Supports two block schemas:
          - aistudio-app: {"block_label": "doc_title", "block_content": "..."}
          - Legacy: {"type": "text", "text": "..."}
        """
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []

        for block in blocks:
            if not isinstance(block, dict):
                continue

            block_type = block.get(
                "block_label",
                block.get("type", block.get("category", "text")),
            )
            content = block.get(
                "block_content",
                block.get("text", block.get("content", "")),
            )
            if not content:
                continue

            if block_type in ("title", "heading", "doc_title"):
                text_parts.append(f"# {content}")
            elif block_type in ("title2", "section_title"):
                text_parts.append(f"## {content}")
            elif block_type == "table":
                rows = block.get("rows", block.get("table_body", []))
                if isinstance(rows, list) and rows:
                    tbl = []
                    for row in rows:
                        if isinstance(row, list):
                            tbl.append([str(c) for c in row])
                        elif isinstance(row, dict):
                            tbl.append([row.get("text", str(row))])
                    if tbl:
                        tables.append(tbl)
                        text_parts.append("[TABLE]\n" + "\n".join(" | ".join(r) for r in tbl))
                    else:
                        text_parts.append(f"[TABLE]\n{content}")
                else:
                    text_parts.append(f"[TABLE]\n{content}")
            elif block_type in ("figure", "image"):
                text_parts.append(f"[FIGURE]\n{content}" if content else "")
            elif block_type == "formula":
                text_parts.append(f"[FORMULA]\n{content}")
            elif block_type == "list":
                text_parts.append(content)
            else:
                text_parts.append(content)

        return "\n\n".join(p for p in text_parts if p), tables

    async def _parse_image(
        self,
        client,
        base_url: str,
        api_key: str,
        file_path: Path,
    ) -> ParsedDocument:
        """Parse a single image file."""
        logger.info("PaddleOCR parsing image: %s", file_path.name)
        img_bytes = file_path.read_bytes()
        blocks = await self._call_layout_api(client, base_url, api_key, img_bytes)
        text, tables = self._blocks_to_text(blocks)

        images = None
        if settings.RAG_INDEX_IMAGE_CAPTIONS:
            images = [
                {
                    "name": file_path.name,
                    "data": base64.b64encode(img_bytes).decode("ascii"),
                }
            ]

        return ParsedDocument(
            file_type=file_path.suffix.lstrip("."),
            pages=[
                PageContent(
                    page_num=1,
                    text=text,
                    tables=tables or None,
                    images=images,
                    metadata={"parser": "paddle", "ocr": True},
                )
            ],
            metadata={"parser": "paddle", "has_ocr": True},
            total_pages=1,
        )

    async def _parse_pdf(
        self,
        client,
        base_url: str,
        api_key: str,
        file_path: Path,
    ) -> ParsedDocument:
        """Parse a PDF — render each page to image and send to API."""
        import fitz

        logger.info("PaddleOCR parsing PDF (per-page): %s", file_path.name)
        scale = settings.OCR_DPI / 72.0

        page_images: list[tuple[int, bytes]] = []
        with fitz.open(str(file_path)) as doc:
            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale))  # type: ignore[union-attr]
                img_bytes = pix.tobytes("png")
                page_images.append((page_num + 1, img_bytes))

        # Process pages with bounded concurrency
        semaphore = asyncio.Semaphore(_MAX_CONCURRENT_PAGES)

        async def _process_page(page_num: int, img_bytes: bytes) -> PageContent:
            async with semaphore:
                blocks = await self._call_layout_api(client, base_url, api_key, img_bytes)
                text, tables = self._blocks_to_text(blocks)
                images = None
                if settings.RAG_INDEX_IMAGE_CAPTIONS:
                    images = [
                        {
                            "name": f"page-{page_num}.png",
                            "data": base64.b64encode(img_bytes).decode("ascii"),
                        }
                    ]
                return PageContent(
                    page_num=page_num,
                    text=text,
                    tables=tables or None,
                    images=images,
                    metadata={"parser": "paddle", "ocr": True},
                )

        tasks = [_process_page(pn, img) for pn, img in page_images]
        pages = await asyncio.gather(*tasks)

        return ParsedDocument(
            file_type="pdf",
            pages=list(pages),
            metadata={"parser": "paddle", "has_ocr": True},
            total_pages=len(pages),
        )

    async def _parse_pptx(
        self,
        client,
        base_url: str,
        api_key: str,
        file_path: Path,
    ) -> ParsedDocument:
        """Parse a PPTX — extract text + OCR image shapes."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        logger.info("PaddleOCR parsing PPTX: %s", file_path.name)
        prs = Presentation(str(file_path))
        pages: list[PageContent] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts: list[str] = []
            tables: list[list[list[str]]] = []
            slide_images: list[dict[str, str]] = []

            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if text:
                    text_parts.append(text)

                if getattr(shape, "has_table", False):
                    tbl = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows  # type: ignore[attr-defined]
                    ]
                    tables.append(tbl)
                    text_parts.append("[TABLE]\n" + "\n".join(" | ".join(r) for r in tbl))

                # OCR embedded image shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob  # type: ignore[attr-defined]
                        if settings.RAG_INDEX_IMAGE_CAPTIONS:
                            image_ext = shape.image.ext or "png"  # type: ignore[attr-defined]
                            slide_images.append(
                                {
                                    "name": f"slide-{slide_num}-image.{image_ext}",
                                    "data": base64.b64encode(img_blob).decode("ascii"),
                                }
                            )
                        blocks = await self._call_layout_api(client, base_url, api_key, img_blob)
                        ocr_text, _ = self._blocks_to_text(blocks)
                        if ocr_text:
                            text_parts.append(f"[IMAGE OCR]\n{ocr_text}")
                    except Exception as exc:
                        logger.debug("Image shape OCR skipped: %s", exc)

            # Speaker notes
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"[NOTES]\n{notes}")

            pages.append(
                PageContent(
                    page_num=slide_num,
                    text="\n\n".join(text_parts),
                    tables=tables or None,
                    images=slide_images or None,
                    metadata={"parser": "paddle", "ocr": True, "has_notes": bool(notes)},
                )
            )

        return ParsedDocument(
            file_type="pptx",
            pages=pages,
            metadata={
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "parser": "paddle",
                "has_ocr": True,
            },
            total_pages=len(pages),
        )
