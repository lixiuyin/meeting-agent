"""Local fast parser — pure-local text extraction via PyMuPDF, python-pptx, etc.

No OCR, no cloud calls.  Used as both primary (for clean text docs) and
last-resort fallback (when all cloud APIs are down).
"""

import base64
import logging
from pathlib import Path

from src.core.config import settings

from ..types import PageContent, ParsedDocument

logger = logging.getLogger(__name__)

# Image extensions that should use the image-specific parsing path.
_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"})


class LocalFastParser:
    """Deterministic, free, local-only text extraction."""

    def parse(self, file_path: Path) -> ParsedDocument:
        """Parse a document using local libraries only."""
        suffix = file_path.suffix.lower()

        if suffix == ".pdf":
            return self._parse_pdf(file_path)
        if suffix == ".pptx":
            return self._parse_pptx(file_path)
        if suffix == ".docx":
            return self._parse_docx(file_path)
        if suffix == ".xlsx":
            return self._parse_xlsx(file_path)
        if suffix in _IMAGE_EXTS:
            return self._parse_image(file_path)

        # Generic fallback — try PyMuPDF first, then plain read
        return self._parse_generic(file_path)

    def _parse_image(self, file_path: Path) -> ParsedDocument:
        """Local OCR fallback for standalone image files.

        Tries pytesseract → easyocr → graceful placeholder.  This ensures
        users always get a usable result even when all cloud OCR APIs are
        down.
        """
        logger.info("Local-fast parsing image: %s", file_path.name)

        text = self._try_tesseract(file_path)
        ocr_engine = "pytesseract"

        if not text:
            text = self._try_easyocr(file_path)
            ocr_engine = "easyocr"

        if not text:
            # No local OCR available — try vision service as a last resort
            # so the pipeline can still archive meaningful text.
            ocr_engine = "none"
            text = ""
            logger.warning(
                "No local OCR available for %s — trying vision fallback",
                file_path.name,
            )
            try:
                import asyncio as _asyncio

                from ...vision import caption_image, transcribe_text_bearing_image

                async def _process():
                    caption, ocr = await _asyncio.gather(
                        caption_image(str(file_path)),
                        transcribe_text_bearing_image(str(file_path)),
                    )
                    return caption, ocr

                caption, ocr = _asyncio.run(_process())

                parts = [p for p in (caption, ocr) if p]
                if parts:
                    text = "\n\n".join(parts)
                    ocr_engine = "vision"
                    logger.info(
                        "Vision fallback extracted %d chars from %s",
                        len(text),
                        file_path.name,
                    )
            except Exception:
                logger.warning(
                    "Vision fallback also failed for %s",
                    file_path.name,
                    exc_info=True,
                )

        # Encode image as a base64 asset so downstream processors can
        # generate thumbnails / captions.
        image_b64: str | None = None
        try:
            image_b64 = base64.b64encode(file_path.read_bytes()).decode("ascii")
        except Exception:
            logger.debug("Failed to read image bytes for %s", file_path.name, exc_info=True)

        metadata: dict = {
            "parser": "local",
            "has_ocr": ocr_engine != "none",
            "ocr_engine": ocr_engine,
        }
        if ocr_engine == "none":
            metadata["local_ocr_unavailable"] = True

        return ParsedDocument(
            file_type=file_path.suffix.lstrip("."),
            pages=[
                PageContent(
                    page_num=1,
                    text=text,
                    images=[{"name": file_path.name, "data": image_b64}] if image_b64 else None,
                    metadata=metadata,
                )
            ],
            metadata=metadata,
            total_pages=1,
        )

    @staticmethod
    def _try_tesseract(file_path: Path) -> str:
        """Attempt OCR with pytesseract.  Returns empty string if unavailable."""
        try:
            import pytesseract  # type: ignore[import-untyped]
            from PIL import Image

            img = Image.open(str(file_path))
            text = pytesseract.image_to_string(img).strip()
            if text:
                logger.debug("pytesseract extracted %d chars from %s", len(text), file_path.name)
            return text
        except ImportError:
            logger.debug("pytesseract not installed, skipping")
            return ""
        except Exception:
            logger.debug("pytesseract failed for %s", file_path.name, exc_info=True)
            return ""

    @staticmethod
    def _try_easyocr(file_path: Path) -> str:
        """Attempt OCR with easyocr (CPU).  Returns empty string if unavailable."""
        try:
            import easyocr  # type: ignore[import-untyped]

            reader = easyocr.Reader(["en", "ch_sim"], verbose=False)
            results = reader.readtext(str(file_path))
            text = "\n".join(item[1] for item in results if item[1]).strip()
            if text:
                logger.debug("easyocr extracted %d chars from %s", len(text), file_path.name)
            return text
        except ImportError:
            logger.debug("easyocr not installed, skipping")
            return ""
        except Exception:
            logger.debug("easyocr failed for %s", file_path.name, exc_info=True)
            return ""

    def _parse_pdf(self, file_path: Path) -> ParsedDocument:
        """Extract text from PDF using PyMuPDF."""
        import pymupdf as fitz

        logger.info("Local-fast parsing PDF: %s", file_path.name)
        pages: list[PageContent] = []

        with fitz.open(str(file_path)) as doc:
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text("text").strip()  # type: ignore[union-attr]
                images: list[dict] = []
                try:
                    for img_idx, image_info in enumerate(page.get_images(full=True)):  # type: ignore[union-attr]
                        if len(images) >= settings.PARSER_MAX_IMAGES_PER_PAGE:
                            break
                        xref = image_info[0]
                        extracted = doc.extract_image(xref)
                        image_bytes = extracted.get("image")
                        if not isinstance(image_bytes, bytes) or not image_bytes:
                            continue
                        if len(image_bytes) > settings.PARSER_MAX_IMAGE_BYTES:
                            continue
                        ext = extracted.get("ext", "png")
                        images.append(
                            {
                                "name": f"page-{page_num + 1}-img-{img_idx}.{ext}",
                                "data": base64.b64encode(image_bytes).decode("ascii"),
                            }
                        )
                except Exception:
                    logger.debug("Local PDF image extraction skipped", exc_info=True)
                if settings.RAG_INDEX_IMAGE_CAPTIONS and not text and not images:
                    # Some scanned PDFs expose no extractable text and no embedded
                    # image objects. Render the page as a fallback image so
                    # downstream multimodal indexing can still proceed.
                    try:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # type: ignore[union-attr]
                        rendered_bytes = pix.tobytes("png")
                        if len(rendered_bytes) > settings.PARSER_MAX_IMAGE_BYTES:
                            logger.debug(
                                "Local PDF render fallback skipped due to size limit (%d bytes)",
                                len(rendered_bytes),
                            )
                            raise ValueError("Rendered page image exceeds size limit")
                        images.append(
                            {
                                "name": f"page-{page_num + 1}.png",
                                "data": base64.b64encode(rendered_bytes).decode("ascii"),
                            }
                        )
                    except Exception:
                        logger.debug("Local PDF page render fallback skipped", exc_info=True)
                pages.append(
                    PageContent(
                        page_num=page_num + 1,
                        text=text,
                        images=images or None,
                        metadata={"ocr": False},
                    )
                )

        return ParsedDocument(
            file_type="pdf",
            pages=pages,
            metadata={"parser": "local", "has_ocr": False},
            total_pages=len(pages),
        )

    def _parse_pptx(self, file_path: Path) -> ParsedDocument:
        """Extract text from PPTX slides via python-pptx."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        logger.info("Local-fast parsing PPTX: %s", file_path.name)
        prs = Presentation(str(file_path))
        pages: list[PageContent] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts: list[str] = []
            tables: list[list[list[str]]] = []
            images: list[dict] = []

            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if text:
                    text_parts.append(text)

                if getattr(shape, "has_table", False):
                    tbl = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows  # type: ignore[attr-defined]
                    ]
                    tables.append(tbl)
                    text_parts.append("[TABLE]\n" + "\n".join(" | ".join(r) for r in tbl))
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if len(images) >= settings.PARSER_MAX_IMAGES_PER_PAGE:
                        break
                    try:
                        image_bytes = shape.image.blob  # type: ignore[attr-defined]
                        if len(image_bytes) > settings.PARSER_MAX_IMAGE_BYTES:
                            continue
                        image_ext = shape.image.ext or "png"  # type: ignore[attr-defined]
                        images.append(
                            {
                                "name": f"slide-{slide_num}-image.{image_ext}",
                                "data": base64.b64encode(image_bytes).decode("ascii"),
                            }
                        )
                    except Exception:
                        logger.debug("PPTX image extraction skipped", exc_info=True)

            # Speaker notes
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"[NOTES]\n{notes}")

            pages.append(
                PageContent(
                    page_num=slide_num,
                    text="\n\n".join(text_parts),
                    tables=tables or None,
                    images=images or None,
                    metadata={"ocr": False, "has_notes": bool(notes)},
                )
            )

        return ParsedDocument(
            file_type="pptx",
            pages=pages,
            metadata={
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "parser": "local",
                "has_ocr": False,
            },
            total_pages=len(pages),
        )

    def _parse_docx(self, file_path: Path) -> ParsedDocument:
        """Extract text from DOCX via python-docx."""
        from docx import Document as DocxDocument  # type: ignore[import-not-found]

        logger.info("Local-fast parsing DOCX: %s", file_path.name)
        doc = DocxDocument(str(file_path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        return ParsedDocument(
            file_type="docx",
            pages=[
                PageContent(
                    page_num=1,
                    text="\n\n".join(paragraphs),
                    metadata={"ocr": False},
                )
            ],
            metadata={"parser": "local", "has_ocr": False},
            total_pages=1,
        )

    def _parse_xlsx(self, file_path: Path) -> ParsedDocument:
        """Extract text from XLSX via openpyxl."""
        from openpyxl import load_workbook  # type: ignore[import-not-found]

        logger.info("Local-fast parsing XLSX: %s", file_path.name)
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        text_parts: list[str] = []

        for sheet in wb.worksheets:
            sheet_text: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    sheet_text.append(" | ".join(cells))
            if sheet_text:
                text_parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(sheet_text))

        wb.close()

        return ParsedDocument(
            file_type="xlsx",
            pages=[
                PageContent(
                    page_num=1,
                    text="\n\n".join(text_parts),
                    metadata={"ocr": False},
                )
            ],
            metadata={"parser": "local", "has_ocr": False},
            total_pages=1,
        )

    def _parse_generic(self, file_path: Path) -> ParsedDocument:
        """Last-resort: try PyMuPDF, then plain text read."""
        suffix = file_path.suffix.lower()
        text = ""

        # Try PyMuPDF for anything it can open
        try:
            import pymupdf as fitz

            with fitz.open(str(file_path)) as doc:
                for page in doc:
                    text += page.get_text("text") + "\n"  # type: ignore[union-attr]
        except Exception:
            pass  # PyMuPDF text extraction best-effort; caller falls back to other parsers

        if not text.strip():
            try:
                text = file_path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                text = ""

        return ParsedDocument(
            file_type=suffix.lstrip("."),
            pages=[
                PageContent(
                    page_num=1,
                    text=text.strip(),
                    metadata={"ocr": False},
                )
            ],
            metadata={"parser": "local", "has_ocr": False},
            total_pages=1,
        )
