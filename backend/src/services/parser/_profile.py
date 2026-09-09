"""Document profiling — fast local pre-analysis to drive parser routing.

Uses only PyMuPDF (fitz) for PDFs and python-pptx for PPTX — no cloud calls,
no heavy ML deps.  Always returns a profile; never raises (falls back to a
generic "unknown" profile on error so the router still picks something).
"""

import contextlib
import logging
import os
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"})
_TEXT_EXTS = frozenset(
    {".txt", ".md", ".markdown", ".html", ".htm", ".json", ".xml", ".rtf", ".csv"}
)


@dataclass(frozen=True)
class DocumentProfile:
    """Immutable summary of a document's characteristics for routing."""

    suffix: str
    page_count: int
    avg_chars_per_page: float
    image_ratio: float  # 0.0-1.0, fraction of pages with significant images
    has_text_layer: bool
    is_likely_scanned: bool
    size_bytes: int


def profile_document(file_path: Path) -> DocumentProfile:
    """Analyse a document and return a profile for router selection.

    Never raises — returns a conservative "unknown" profile on any error.
    """
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    size_bytes: int = 0
    with contextlib.suppress(OSError):
        size_bytes = os.path.getsize(file_path)

    # Single image files — no pages to analyse
    if suffix in _IMAGE_EXTS:
        return DocumentProfile(
            suffix=suffix,
            page_count=1,
            avg_chars_per_page=0.0,
            image_ratio=1.0,
            has_text_layer=False,
            is_likely_scanned=False,
            size_bytes=size_bytes,
        )

    if suffix == ".pdf":
        return _profile_pdf(file_path, size_bytes)

    if suffix == ".pptx":
        return _profile_pptx(file_path, size_bytes)

    # DOCX / XLSX / other formats — treat as text-present, no images
    return DocumentProfile(
        suffix=suffix,
        page_count=1,
        avg_chars_per_page=0.0,
        image_ratio=0.0,
        has_text_layer=True,
        is_likely_scanned=False,
        size_bytes=size_bytes,
    )


def _profile_pdf(file_path: Path, size_bytes: int) -> DocumentProfile:
    """Profile a PDF using PyMuPDF (fitz)."""
    try:
        import pymupdf as fitz

        with fitz.open(str(file_path)) as doc:
            page_count = len(doc)
            if page_count == 0:
                return _unknown_profile(".pdf", size_bytes)

            total_chars = 0
            image_pages = 0

            for page in doc:
                raw_text = page.get_text("text")  # type: ignore[union-attr]
                text = raw_text if isinstance(raw_text, str) else ""
                total_chars += len(text.strip())
                # Count images on this page
                img_list = page.get_images(full=True)  # type: ignore[union-attr]
                if img_list:
                    image_pages += 1

            avg_chars = total_chars / page_count
            image_ratio = image_pages / page_count

            return DocumentProfile(
                suffix=".pdf",
                page_count=page_count,
                avg_chars_per_page=avg_chars,
                image_ratio=image_ratio,
                has_text_layer=avg_chars > 0,
                is_likely_scanned=avg_chars < 50,
                size_bytes=size_bytes,
            )
    except Exception:
        logger.debug("PDF profiling failed for %s", file_path.name, exc_info=True)
        return _unknown_profile(".pdf", size_bytes)


def _profile_pptx(file_path: Path, size_bytes: int) -> DocumentProfile:
    """Profile a PPTX using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(file_path))
        slide_count = len(prs.slides)
        if slide_count == 0:
            return _unknown_profile(".pptx", size_bytes)

        total_chars = 0
        image_heavy_slides = 0

        for slide in prs.slides:
            slide_chars = 0
            has_images = False
            for shape in slide.shapes:
                raw_text = getattr(shape, "text", "")
                text = raw_text.strip() if isinstance(raw_text, str) else ""
                slide_chars += len(text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_images = True
            total_chars += slide_chars
            # A slide is "image-heavy" if it has images and very little text
            if has_images and slide_chars < 50:
                image_heavy_slides += 1

        avg_chars = total_chars / slide_count
        image_ratio = image_heavy_slides / slide_count

        return DocumentProfile(
            suffix=".pptx",
            page_count=slide_count,
            avg_chars_per_page=avg_chars,
            image_ratio=image_ratio,
            has_text_layer=avg_chars > 0,
            is_likely_scanned=False,
            size_bytes=size_bytes,
        )
    except Exception:
        logger.debug("PPTX profiling failed for %s", file_path.name, exc_info=True)
        return _unknown_profile(".pptx", size_bytes)


def _unknown_profile(suffix: str, size_bytes: int) -> DocumentProfile:
    """Fallback profile when analysis fails — triggers cloud parsing."""
    return DocumentProfile(
        suffix=suffix,
        page_count=0,
        avg_chars_per_page=0.0,
        image_ratio=0.0,
        has_text_layer=False,
        is_likely_scanned=False,
        size_bytes=size_bytes,
    )
