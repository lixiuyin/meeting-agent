"""Legacy format converters — PPT, DOC, XLS to modern formats."""

import contextlib
import logging
import shutil
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

# Probe once at module load — LibreOffice is optional for legacy formats
_LIBREOFFICE_AVAILABLE = shutil.which("soffice") is not None
if not _LIBREOFFICE_AVAILABLE:
    logger.warning(
        "LibreOffice (soffice) not found. Legacy .ppt/.doc/.xls conversion unavailable. "
        "Install with: apt install libreoffice (Debian/Ubuntu) or brew install libreoffice (macOS)"
    )


class LibreOfficeMissingError(RuntimeError):
    """Raised when LibreOffice is required but not installed."""


def _convert_via_libreoffice(src_path: Path, target_format: str) -> Path:
    """Generic LibreOffice conversion. Raises LibreOfficeMissingError if soffice missing.

    Uses a dedicated user profile directory per invocation so concurrent
    conversions don't collide (LibreOffice only allows one instance per profile).
    """
    import subprocess

    if not _LIBREOFFICE_AVAILABLE:
        raise LibreOfficeMissingError(
            f"LibreOffice is required to convert {src_path.suffix} files but is not installed. "
            "Install with: apt install libreoffice (Debian/Ubuntu) "
            "or brew install libreoffice (macOS)"
        )

    tmp_dir = tempfile.mkdtemp()
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        # file:// URL works on all platforms; Path.as_uri() produces
        # file:/// on Unix and file:///C:/ on Windows.
        profile_url = Path(profile_dir).as_uri()
        cmd = [
            "soffice",
            f"-env:UserInstallation={profile_url}",
            "--headless",
            "--norestore",
            "--nologo",
            "--convert-to",
            target_format,
            "--outdir",
            tmp_dir,
            str(src_path),
        ]
        subprocess.run(cmd, check=True, capture_output=True, timeout=120)

        output_name = src_path.stem + f".{target_format}"
        output_path = Path(tmp_dir) / output_name
        if not output_path.exists():
            raise FileNotFoundError(
                f"LibreOffice conversion failed: {output_name} not found in {tmp_dir}"
            )
        return output_path
    except Exception:
        # Clean up temp dir on failure; on success the caller is responsible
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


def _convert_ppt_to_pptx(ppt_path: Path) -> Path:
    """Convert legacy .ppt to .pptx using LibreOffice."""
    return _convert_via_libreoffice(ppt_path, "pptx")


def _convert_doc_to_docx(doc_path: Path) -> Path:
    """Convert legacy .doc to .docx using LibreOffice."""
    return _convert_via_libreoffice(doc_path, "docx")


def _convert_xls_to_xlsx(xls_path: Path) -> Path:
    """Convert legacy .xls to .xlsx using LibreOffice."""
    return _convert_via_libreoffice(xls_path, "xlsx")


def _convert_pptx_to_pdf(pptx_path: Path) -> Path:
    """Convert .pptx to .pdf.

    Primary: LibreOffice (full fidelity, preserves layout/images).
    Fallback: python-pptx text extraction + PyMuPDF PDF generation.
    """
    try:
        return _convert_via_libreoffice(pptx_path, "pdf")
    except (LibreOfficeMissingError, OSError, Exception) as exc:
        if not isinstance(exc, (LibreOfficeMissingError, OSError)):
            # Re-raise unexpected errors (e.g. PermissionError)
            # but continue for subprocess failures (CalledProcessError, TimeoutExpired)
            import subprocess

            if not isinstance(exc, subprocess.CalledProcessError | subprocess.TimeoutExpired):
                raise
        logger.info(
            "LibreOffice unavailable for PPTX→PDF (%s). "
            "Falling back to python-pptx + PyMuPDF conversion.",
            exc,
        )
        return _pptx_to_pdf_via_extract(pptx_path)


def _convert_ppt_to_pdf(ppt_path: Path) -> Path:
    """Convert legacy .ppt directly to .pdf.

    Primary: LibreOffice direct .ppt→.pdf.
    Fallback: LibreOffice .ppt→.pptx then python-pptx + PyMuPDF.
    Last resort: python-pptx + PyMuPDF (only if .ppt can be read directly).
    """
    try:
        return _convert_via_libreoffice(ppt_path, "pdf")
    except (LibreOfficeMissingError, OSError) as exc:
        logger.info("LibreOffice unavailable for .ppt→.pdf (%s). Trying .ppt→.pptx fallback.", exc)
    import subprocess

    pptx_path: Path | None = None
    try:
        pptx_path = _convert_via_libreoffice(ppt_path, "pptx")
    except (LibreOfficeMissingError, OSError, subprocess.CalledProcessError):
        logger.info("LibreOffice .ppt→.pptx also failed. Attempting direct python-pptx read.")
        return _pptx_to_pdf_via_extract(ppt_path)
    try:
        return _convert_via_libreoffice(pptx_path, "pdf")
    except (LibreOfficeMissingError, OSError, subprocess.CalledProcessError):
        return _pptx_to_pdf_via_extract(pptx_path)
    finally:
        with contextlib.suppress(OSError):
            pptx_path.unlink(missing_ok=True)
        with contextlib.suppress(OSError):
            pptx_path.parent.rmdir()


def _pptx_to_pdf_via_extract(pptx_path: Path) -> Path:
    """Fallback PPTX→PDF using python-pptx text extraction + PyMuPDF.

    Extracts text, tables, and speaker notes from each slide and writes
    a simple text-based PDF. Visual fidelity is reduced but the content
    is fully preserved for the cascade OCR/parsing pipeline.
    """
    import pymupdf as fitz
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    doc = fitz.open()

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts: list[str] = []

        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                text_parts.append(text)

            if getattr(shape, "has_table", False):
                tbl = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows  # type: ignore[attr-defined]
                ]
                text_parts.append("[TABLE]\n" + "\n".join(" | ".join(r) for r in tbl))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_parts.append(f"[NOTES]\n{notes}")

        # Extract images from shapes and embed as full pages
        slide_images: list[bytes] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob  # type: ignore[attr-defined]
                    if len(img_bytes) <= 20 * 1024 * 1024:
                        slide_images.append(img_bytes)
                except Exception:
                    pass  # per-shape image extraction; skip broken shapes

        page = doc.new_page(width=612, height=792)

        # If slide has images, embed the first image as a full-page visual
        if slide_images:
            try:
                img_doc = fitz.open("png", slide_images[0])
                rect = page.rect
                page.insert_image(rect, pixmap=img_doc[0].get_pixmap())
                img_doc.close()
            except Exception:
                pass  # image insertion best-effort; text fallback still works

        # Insert text below or as the main content
        combined_text = f"--- Slide {slide_num} ---\n\n" + "\n\n".join(text_parts)
        if combined_text.strip():
            text_rect = (
                page.rect
                if not slide_images
                else fitz.Rect(0, page.rect.height * 0.6, page.rect.width, page.rect.height)
            )
            page.insert_textbox(
                text_rect,
                combined_text,
                fontsize=10,
                fontname="helv",
            )

    tmp_dir = tempfile.mkdtemp()
    output_path = Path(tmp_dir) / (pptx_path.stem + ".pdf")
    try:
        doc.save(str(output_path))
        doc.close()
        logger.info(
            "Converted PPTX→PDF via python-pptx+PyMuPDF: %d slides → %s",
            len(prs.slides),
            output_path,
        )
        return output_path
    except Exception:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise
